from __future__ import annotations

import csv
import hashlib
import html
import logging
import os
import re
import shutil
import subprocess
import tarfile
import tempfile
import time
import zipfile
from pathlib import Path

from bs4 import BeautifulSoup
from odf import teletype
from odf.opendocument import load as load_odf_document
from pypdf import PdfReader
import docx
from pptx import Presentation
from striprtf.striprtf import rtf_to_text


XUTF8_BINARY = Path(os.getenv("XUTF8_BINARY_PATH", "/app/bin/xutf_8"))
XUTF8_TIMEOUT_SEC = int(os.getenv("XUTF8_TIMEOUT_SEC", "60"))
logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".doc",
    ".docx",
    ".pdf",
    ".py",
    ".txt",
    ".hwp",
    ".hwpx",
    ".odt",
    ".ppt",
    ".pptx",
    ".xls",
    ".xlsx",
    ".csv",
    ".html",
    ".htm",
    ".rtf",
    ".md",
}

TEXT_EXTENSIONS = {
    ".txt",
    ".text",
    ".log",
    ".md",
    ".markdown",
    ".rst",
    ".json",
    ".jsonl",
    ".xml",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".conf",
    ".cfg",
    ".properties",
    ".env",
    ".sql",
    ".graphql",
    ".proto",
    ".csv",
    ".tsv",
    ".html",
    ".htm",
    ".css",
    ".scss",
    ".sass",
    ".less",
    ".py",
    ".pyw",
    ".ipynb",
    ".js",
    ".jsx",
    ".ts",
    ".tsx",
    ".mjs",
    ".cjs",
    ".java",
    ".kt",
    ".kts",
    ".go",
    ".rs",
    ".c",
    ".h",
    ".cpp",
    ".cc",
    ".cxx",
    ".hpp",
    ".cs",
    ".php",
    ".rb",
    ".pl",
    ".pm",
    ".r",
    ".scala",
    ".swift",
    ".dart",
    ".lua",
    ".groovy",
    ".gradle",
    ".sh",
    ".bash",
    ".zsh",
    ".fish",
    ".ps1",
    ".bat",
    ".cmd",
    ".vue",
    ".svelte",
}

SUPPORTED_EXTENSIONS.update(TEXT_EXTENSIONS)

ARCHIVE_EXTENSIONS = {".zip", ".tar", ".tar.gz", ".tgz"}
UPLOAD_EXTENSIONS = SUPPORTED_EXTENSIONS | ARCHIVE_EXTENSIONS


def detect_extension(path_or_name: str | Path) -> str:
    name = Path(path_or_name).name.lower()
    for suffix in (".tar.gz", ".tgz"):
        if name.endswith(suffix):
            return suffix
    return Path(name).suffix.lower()


def is_supported_single_document(path_or_name: str | Path) -> bool:
    return detect_extension(path_or_name) in SUPPORTED_EXTENSIONS


def is_supported_upload(path_or_name: str | Path) -> bool:
    return detect_extension(path_or_name) in UPLOAD_EXTENSIONS


def extract_text(path: str | Path) -> str:
    started = time.perf_counter()
    p = Path(path)
    suffix = detect_extension(p)
    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {suffix or '(none)'}")

    xutf8_text = _extract_with_xutf8(p)
    if xutf8_text:
        text = _normalize_text(xutf8_text)
        logger.info(
            "document extraction completed method=xutf_8 file=%s ext=%s size=%d chars=%d elapsed_ms=%.1f",
            p.name,
            suffix,
            _safe_file_size(p),
            len(text),
            _elapsed_ms(started),
        )
        return text

    logger.info(
        "document extraction fallback method=internal file=%s ext=%s size=%d reason=xutf_8_unavailable_or_empty",
        p.name,
        suffix,
        _safe_file_size(p),
    )

    if suffix == ".pdf":
        reader = PdfReader(str(p))
        parts = []
        for page in reader.pages:
            text = page.extract_text() or ""
            if text.strip():
                parts.append(text)
        return _log_extracted_text(p, suffix, "pypdf", "\n\n".join(parts), started)

    if suffix == ".docx":
        document = docx.Document(str(p))
        parts = [para.text for para in document.paragraphs if (para.text or "").strip()]
        return _log_extracted_text(p, suffix, "python-docx", "\n".join(parts), started)

    if suffix == ".pptx":
        presentation = Presentation(str(p))
        parts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                text = (getattr(shape, "text", "") or "").strip()
                if text:
                    parts.append(text)
        return _log_extracted_text(p, suffix, "python-pptx", "\n".join(parts), started)

    if suffix == ".xlsx":
        import openpyxl

        workbook = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
        parts = []
        try:
            for sheet in workbook.worksheets:
                text = _join_table_rows([list(row) for row in sheet.iter_rows(values_only=True)])
                if text:
                    parts.append(text)
        finally:
            workbook.close()
        return _log_extracted_text(p, suffix, "openpyxl", "\n\n".join(parts), started)

    if suffix == ".xls":
        import xlrd

        book = xlrd.open_workbook(str(p), on_demand=True)
        parts = []
        try:
            for sheet in book.sheets():
                rows = [sheet.row_values(row_idx) for row_idx in range(sheet.nrows)]
                text = _join_table_rows(rows)
                if text:
                    parts.append(text)
        finally:
            book.release_resources()
        return _log_extracted_text(p, suffix, "xlrd", "\n\n".join(parts), started)

    if suffix == ".csv":
        raw = _read_text_file(p)
        return _log_extracted_text(p, suffix, "csv", _join_table_rows(list(csv.reader(raw.splitlines()))), started)

    if suffix == ".tsv":
        raw = _read_text_file(p)
        return _log_extracted_text(p, suffix, "tsv", _join_table_rows(list(csv.reader(raw.splitlines(), delimiter="\t"))), started)

    if suffix in {".html", ".htm"}:
        raw = _read_text_file(p)
        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        return _log_extracted_text(p, suffix, "html", soup.get_text(separator="\n"), started)

    if suffix in {".md", ".markdown"}:
        return _log_extracted_text(p, suffix, "markdown", _strip_markdown(_read_text_file(p)), started)

    if suffix == ".rtf":
        raw = _read_text_file(p)
        return _log_extracted_text(p, suffix, "rtf", rtf_to_text(raw), started)

    if suffix == ".odt":
        document = load_odf_document(str(p))
        return _log_extracted_text(p, suffix, "odf", teletype.extractText(document), started)

    if suffix == ".doc":
        text = _extract_with_cmd(["antiword", str(p)])
        if text:
            return _log_extracted_text(p, suffix, "antiword", text, started)
        text, error = _extract_with_soffice_to_txt(p)
        if text:
            return _log_extracted_text(p, suffix, "soffice", text, started)
        raise ValueError(f".doc 텍스트 추출 실패: {error or 'antiword/libreoffice 사용 불가'}")

    if suffix == ".ppt":
        text = _extract_with_cmd(["catppt", "-b", str(p)])
        if text:
            return _log_extracted_text(p, suffix, "catppt", text, started)
        text, error = _extract_with_soffice_to_txt(p)
        if text:
            return _log_extracted_text(p, suffix, "soffice", text, started)
        raise ValueError(f".ppt 텍스트 추출 실패: {error or 'catppt/libreoffice 사용 불가'}")

    if suffix in {".hwp", ".hwpx"}:
        text, parser_error = _extract_with_hwp_parser(p)
        if text:
            return _log_extracted_text(p, suffix, "hwp_hwpx_parser", text, started)
        text, soffice_error = _extract_with_soffice_to_txt(p)
        if text:
            return _log_extracted_text(p, suffix, "soffice", text, started)
        detail = "; ".join(item for item in (parser_error, soffice_error) if item)
        raise ValueError(f"{suffix} 텍스트 추출 실패: {detail or 'hwp/hwpx 추출기 사용 불가'}")

    if suffix in TEXT_EXTENSIONS:
        return _log_extracted_text(p, suffix, "text", _read_text_file(p), started)

    return _log_extracted_text(p, suffix, "text", _read_text_file(p), started)


def extract_archive_text(
    path: str | Path,
    *,
    max_files: int,
    max_total_bytes: int,
    max_member_bytes: int,
) -> tuple[str, dict[str, object]]:
    p = Path(path)
    suffix = detect_extension(p)
    if suffix not in ARCHIVE_EXTENSIONS:
        raise ValueError(f"지원하지 않는 압축 파일 형식입니다: {suffix or '(none)'}")

    entries = _list_archive_entries(p, suffix)
    selected = []
    skipped_unsupported = 0
    skipped_large = 0
    total_bytes = 0
    for entry in entries:
        name = entry["name"]
        size = int(entry["size"] or 0)
        if not is_supported_single_document(name):
            skipped_unsupported += 1
            continue
        if size > max_member_bytes:
            skipped_large += 1
            continue
        if total_bytes + size > max_total_bytes:
            skipped_large += 1
            continue
        selected.append(entry)
        total_bytes += size
        if len(selected) > max_files:
            raise ValueError(f"압축 파일 내부 문서는 최대 {max_files}개까지 처리할 수 있습니다.")
    if not selected:
        raise ValueError("압축 파일 안에 처리 가능한 문서가 없습니다.")

    parts = []
    extracted = 0
    skipped_failed = 0
    with tempfile.TemporaryDirectory(prefix="xcn_archive_") as tmp:
        tmp_dir = Path(tmp)
        for idx, entry in enumerate(selected, 1):
            out_path = tmp_dir / f"{idx:04d}_{Path(entry['name']).name}"
            try:
                _extract_archive_member(p, suffix, entry["name"], out_path, max_member_bytes=max_member_bytes)
                text = extract_text(out_path)
            except Exception:
                skipped_failed += 1
                continue
            if not text.strip():
                skipped_failed += 1
                continue
            extracted += 1
            parts.append(f"[파일: {entry['name']}]\n{text}")
    if not parts:
        raise ValueError("압축 파일 내부 문서에서 추출된 텍스트가 없습니다.")
    return _normalize_text("\n\n".join(parts)), {
        "archive_file_count": extracted,
        "archive_selected_count": len(selected),
        "archive_uncompressed_bytes": total_bytes,
        "archive_skipped_unsupported": skipped_unsupported,
        "archive_skipped_large": skipped_large,
        "archive_skipped_failed": skipped_failed,
    }


def extract_archive_documents(
    path: str | Path,
    *,
    max_files: int,
    max_total_bytes: int,
    max_member_bytes: int,
) -> tuple[list[dict[str, object]], dict[str, object]]:
    p = Path(path)
    suffix = detect_extension(p)
    if suffix not in ARCHIVE_EXTENSIONS:
        raise ValueError(f"지원하지 않는 압축 파일 형식입니다: {suffix or '(none)'}")

    entries = _list_archive_entries(p, suffix)
    selected = []
    skipped_unsupported = 0
    skipped_large = 0
    total_bytes = 0
    for entry in entries:
        name = str(entry["name"])
        size = int(entry["size"] or 0)
        if not is_supported_single_document(name):
            skipped_unsupported += 1
            continue
        if size > max_member_bytes:
            skipped_large += 1
            continue
        if total_bytes + size > max_total_bytes:
            skipped_large += 1
            continue
        selected.append(entry)
        total_bytes += size
        if len(selected) > max_files:
            raise ValueError(f"압축 파일 내부 문서는 최대 {max_files}개까지 처리할 수 있습니다.")
    if not selected:
        raise ValueError("압축 파일 안에 처리 가능한 문서가 없습니다.")

    documents: list[dict[str, object]] = []
    skipped_failed = 0
    with tempfile.TemporaryDirectory(prefix="xcn_archive_") as tmp:
        tmp_dir = Path(tmp)
        for idx, entry in enumerate(selected, 1):
            member_name = str(entry["name"])
            out_path = tmp_dir / f"{idx:04d}_{Path(member_name).name}"
            try:
                _extract_archive_member(p, suffix, member_name, out_path, max_member_bytes=max_member_bytes)
                text = extract_text(out_path)
            except Exception:
                skipped_failed += 1
                continue
            if not text.strip():
                skipped_failed += 1
                continue
            checksum_sha256 = _sha256_file(out_path)
            documents.append(
                {
                    "text": _normalize_text(text),
                    "file_name": Path(member_name).name,
                    "file_ext": detect_extension(member_name),
                    "file_size": int(entry["size"] or out_path.stat().st_size),
                    "file_checksum_sha256": checksum_sha256,
                    "archive_member_name": member_name,
                }
            )
    if not documents:
        raise ValueError("압축 파일 내부 문서에서 추출된 텍스트가 없습니다.")
    return documents, {
        "archive_file_count": len(documents),
        "archive_selected_count": len(selected),
        "archive_uncompressed_bytes": total_bytes,
        "archive_skipped_unsupported": skipped_unsupported,
        "archive_skipped_large": skipped_large,
        "archive_skipped_failed": skipped_failed,
    }


def _list_archive_entries(path: Path, suffix: str) -> list[dict[str, object]]:
    if suffix == ".zip":
        with zipfile.ZipFile(path) as archive:
            return [
                {"name": item.filename, "size": int(item.file_size or 0)}
                for item in archive.infolist()
                if not item.is_dir() and not _unsafe_archive_name(item.filename)
            ]
    if suffix in {".tar", ".tar.gz", ".tgz"}:
        mode = "r:gz" if suffix in {".tar.gz", ".tgz"} else "r:"
        with tarfile.open(path, mode) as archive:
            return [
                {"name": item.name, "size": int(item.size or 0)}
                for item in archive.getmembers()
                if item.isfile() and not _unsafe_archive_name(item.name)
            ]
    return []


def _extract_archive_member(path: Path, suffix: str, member_name: str, out_path: Path, *, max_member_bytes: int) -> None:
    if suffix == ".zip":
        with zipfile.ZipFile(path) as archive:
            with archive.open(member_name) as src, out_path.open("wb") as dst:
                _copy_limited(src, dst, max_member_bytes)
        return
    mode = "r:gz" if suffix in {".tar.gz", ".tgz"} else "r:"
    with tarfile.open(path, mode) as archive:
        member = archive.getmember(member_name)
        src = archive.extractfile(member)
        if src is None:
            raise ValueError(f"압축 파일 내부 문서를 읽을 수 없습니다: {member_name}")
        with src, out_path.open("wb") as dst:
            _copy_limited(src, dst, max_member_bytes)


def _copy_limited(src, dst, max_bytes: int) -> None:
    total = 0
    while True:
        chunk = src.read(1024 * 1024)
        if not chunk:
            break
        total += len(chunk)
        if total > max_bytes:
            raise ValueError("압축 파일 내부 문서가 허용 크기를 초과했습니다.")
        dst.write(chunk)


def _sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _unsafe_archive_name(name: str) -> bool:
    value = str(name or "").replace("\\", "/")
    return value.startswith("/") or "/../" in f"/{value}" or value in {"", ".", ".."}


def _join_table_rows(rows: list[list[object]]) -> str:
    lines = []
    for row in rows:
        values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
        if values:
            lines.append(" ".join(values))
    return "\n".join(lines)


def _extract_with_cmd(cmd: list[str]) -> str:
    if not shutil.which(cmd[0]):
        return ""
    try:
        proc = subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    except Exception:
        return ""
    return (proc.stdout or "").strip()


def _extract_with_xutf8(path: Path) -> str:
    started = time.perf_counter()
    binary = XUTF8_BINARY
    if not binary.exists():
        logger.warning("xutf_8 unavailable path=%s file=%s", binary, path.name)
        return ""
    try:
        mode = binary.stat().st_mode
        if not mode & 0o100:
            binary.chmod(mode | 0o100)
    except Exception:
        logger.warning("xutf_8 chmod failed path=%s file=%s", binary, path.name, exc_info=True)
        return ""

    commands = (
        ([str(binary), str(path)], None),
        ([str(binary)], path),
    )
    for command, stdin_path in commands:
        try:
            if stdin_path is None:
                proc = subprocess.run(
                    command,
                    check=False,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    timeout=max(1, int(XUTF8_TIMEOUT_SEC)),
                )
            else:
                with stdin_path.open("rb") as source:
                    proc = subprocess.run(
                        command,
                        stdin=source,
                        check=False,
                        stdout=subprocess.PIPE,
                        stderr=subprocess.PIPE,
                        timeout=max(1, int(XUTF8_TIMEOUT_SEC)),
                    )
        except subprocess.TimeoutExpired:
            logger.warning(
                "xutf_8 timeout file=%s timeout_sec=%d elapsed_ms=%.1f",
                path.name,
                XUTF8_TIMEOUT_SEC,
                _elapsed_ms(started),
            )
            continue
        except OSError as exc:
            logger.warning("xutf_8 execution failed file=%s error=%s", path.name, exc)
            continue
        text = proc.stdout.decode("utf-8", errors="replace").strip()
        if proc.returncode == 0 and text:
            return text
        logger.info(
            "xutf_8 returned no usable text file=%s returncode=%s stdout_chars=%d stderr=%s elapsed_ms=%.1f",
            path.name,
            proc.returncode,
            len(text),
            proc.stderr.decode("utf-8", errors="replace").strip()[:300],
            _elapsed_ms(started),
        )
    return ""


def _log_extracted_text(path: Path, suffix: str, method: str, raw_text: str, started: float) -> str:
    text = _normalize_text(raw_text)
    logger.info(
        "document extraction completed method=%s file=%s ext=%s size=%d chars=%d elapsed_ms=%.1f",
        method,
        path.name,
        suffix,
        _safe_file_size(path),
        len(text),
        _elapsed_ms(started),
    )
    return text


def _safe_file_size(path: Path) -> int:
    try:
        return int(path.stat().st_size)
    except OSError:
        return -1


def _elapsed_ms(started: float) -> float:
    return (time.perf_counter() - started) * 1000.0


def _extract_with_soffice_to_txt(path: Path) -> tuple[str, str]:
    if not shutil.which("soffice"):
        return "", "soffice not found"
    try:
        proc = subprocess.run(
            ["soffice", "--headless", "--convert-to", "txt:Text", "--outdir", str(path.parent), str(path)],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )
    except Exception as exc:
        return "", f"soffice convert failed: {exc}"
    txt_path = path.with_suffix(".txt")
    if not txt_path.exists():
        detail = (proc.stderr or proc.stdout or "txt output not created").strip()
        return "", f"soffice no output: {detail}"
    try:
        return txt_path.read_text(encoding="utf-8", errors="ignore").strip(), ""
    finally:
        try:
            txt_path.unlink()
        except Exception:
            pass


def _extract_with_hwp_parser(path: Path) -> tuple[str, str]:
    try:
        from hwp_hwpx_parser import Reader
    except Exception as exc:
        return "", f"hwp_hwpx_parser import failed: {exc}"
    try:
        with Reader(str(path)) as reader:
            return str(getattr(reader, "text", "") or "").strip(), ""
    except Exception as exc:
        return "", f"hwp_hwpx_parser read failed: {exc}"


def _read_text_file(path: Path) -> str:
    raw = path.read_bytes()
    if b"\x00" in raw[:4096]:
        raise ValueError("바이너리 파일로 판단되어 텍스트 추출을 중단했습니다.")
    for encoding in ("utf-8-sig", "utf-8", "cp949", "euc-kr"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="ignore")


def _normalize_text(text: str) -> str:
    return "\n".join(line.strip() for line in str(text or "").splitlines() if line.strip()).strip()


def _strip_markdown(text: str) -> str:
    text = html.unescape(str(text or ""))
    text = re.sub(r"`{1,3}.*?`{1,3}", " ", text, flags=re.DOTALL)
    text = re.sub(r"!\[[^\]]*]\([^)]+\)", " ", text)
    text = re.sub(r"\[([^]]+)]\([^)]+\)", r"\1", text)
    text = re.sub(r"(^|\s)[#>*_~-]+", " ", text, flags=re.MULTILINE)
    return text
